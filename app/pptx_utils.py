from typing import List, Dict, Any, Optional
from pptx import Presentation
from pptx.util import Inches
import io
import random

def collect_template_images(prs: Presentation):
    images = []
    for slide in prs.slides:
        for shape in slide.shapes:
            try:
                if getattr(shape, "shape_type", None) == 13:  # 13 = PICTURE
                    image = shape.image
                    images.append(image.blob)
            except Exception:
                pass
    return images

def add_bullets_to_placeholder(placeholder, bullets: List[str]):
    tf = placeholder.text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = b
        p.level = 0

def add_title_and_content(slide, title: str, bullets: List[str]):
    title_placeholder = None
    body_placeholder = None
    for shape in slide.shapes:
        if getattr(shape, "is_placeholder", False):
            try:
                phf = shape.placeholder_format
            except Exception:
                phf = None
            if phf:
                # 1 => TITLE, 2/7 => BODY/CONTENT (approx)
                if getattr(phf, "type", None) == 1:
                    title_placeholder = shape
                elif getattr(phf, "type", None) in (2, 7):
                    body_placeholder = shape
    if title_placeholder:
        try:
            title_placeholder.text_frame.text = title
        except Exception:
            title_placeholder.text = title
    if body_placeholder:
        add_bullets_to_placeholder(body_placeholder, bullets)

def maybe_add_image(slide, template_images: List[bytes]):
    if not template_images:
        return
    try:
        blob = random.choice(template_images)
        # place near bottom-right; fixed inches keeps behavior consistent
        left = Inches(6)
        top = Inches(4.5)
        slide.shapes.add_picture(io.BytesIO(blob), left, top, width=Inches(2.5))
    except Exception:
        pass

def build_presentation(template_bytes: Optional[bytes], outline: Dict[str, Any]) -> bytes:
    """
    If template_bytes is None, start a blank Presentation() and use its default layouts.
    """
    if template_bytes:
        base = io.BytesIO(template_bytes)
        prs = Presentation(base)
    else:
        prs = Presentation()  # blank presentation (default layouts)

    template_images = collect_template_images(prs)

    # choose first few layouts (fallback to whatever is available)
    layouts = [prs.slide_layouts[i] for i in range(min(5, len(prs.slide_layouts)))]
    if not layouts:
        prs.slides.add_slide(prs.slide_layouts[0])
        return _save(prs)

    for i, slid in enumerate(outline.get("slides", [])):
        layout = layouts[min(i % len(layouts), len(layouts)-1)]
        slide = prs.slides.add_slide(layout)
        title = slid.get("title", f"Slide {i+1}")
        bullets = slid.get("bullets", [])
        add_title_and_content(slide, title, bullets)

        notes = slid.get("notes")
        if notes:
            try:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = notes
            except Exception:
                pass

        # reuse template images occasionally
        if i % 2 == 1:
            maybe_add_image(slide, template_images)

    return _save(prs)

def _save(prs: Presentation) -> bytes:
    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()